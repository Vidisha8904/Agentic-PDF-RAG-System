# chunking with agent, table chunking complete no table is splitted, every table has headers.
# chunking for qa perfect, chunking for normal is done
# table retrieval issues need to refine that
# 25/4

from fastapi import FastAPI, UploadFile, HTTPException
from pydantic import BaseModel
import os
from typing import List
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import ChatOpenAI
from langchain.docstore.document import Document
from dotenv import load_dotenv
from db import add_to_collection, retrieve_from_collection, overwriting_chromadb, delete_pdf_from_collection
from io import BytesIO
import camelot
import json
import requests
import traceback
from fastapi import Depends, HTTPException, status
from fastapi.security import OAuth2PasswordBearer
import requests
from fastapi.middleware.cors import CORSMiddleware
import boto3  # Add boto3 for AWS S3 interaction
import logging
from datetime import datetime
from io import BytesIO
from docx import Document as DocxDocument
from pptx import Presentation
from pdf2image import convert_from_path
import requests
from mistralai import Mistral
from langchain.docstore.document import Document
from io import StringIO
import re
import gc
# import spacy
from typing import List, Dict, Any
from langchain_core.documents import Document
from crewai import Agent, Task, Crew
from crewai.tools import BaseTool
from pydantic import BaseModel, Field, PrivateAttr
import random
import logging
import os
from typing import Dict, Any, List, Tuple  # Added Tuple
from openai import OpenAI
from langchain.tools import BaseTool
from groq import Groq
from typing import ClassVar, Callable, Dict, Any, List, Tuple


# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)
logging.basicConfig(filename="app.log", level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


# Load environment variables
load_dotenv()
os.getenv("OPENAI_API_KEY")
# groq_api_key=os.getenv("groq_api_key")

# Initialize FastAPI app
app = FastAPI(title="PDF RAG Backend")
# cors configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Match React dev server
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# S3 Configuration
S3_BUCKET_NAME = os.getenv("S3_BUCKET_NAME")   
S3_REGION = os.getenv("AWS_REGION") 
aws_access_key = os.getenv("AWS_ACCESS_KEY_ID")
aws_secret_key = os.getenv("AWS_SECRET_ACCESS_KEY")

if not aws_access_key or not aws_secret_key:
    logger.error("AWS credentials not found in environment variables")
    raise ValueError("AWS_ACCESS_KEY_ID and AWS_SECRET_ACCESS_KEY must be set in environment variables")
s3_client = boto3.client(
    "s3",
    region_name=S3_REGION,
    aws_access_key_id=aws_access_key,
    aws_secret_access_key=aws_secret_key
)

# Pydantic models for request/response validation
class URLRequest(BaseModel):
    urls: List[str]

class QueryRequest(BaseModel):
    question: str

class QueryResponse(BaseModel):
    response: str

# Directus configuration
DIRECTUS_URL = "http://directus:8055"

HEADERS = {"Content-Type": "application/json"}

@app.get("/health")
async def health_check():
    return {"status": "healthy"}

@app.options("/query")
async def options_query():
    return {"message": "CORS preflight handled"}

# Load environment variables
load_dotenv()
MISTRAL_API_KEY = os.getenv("MISTRAL_API_KEY")
if not MISTRAL_API_KEY:
    raise ValueError("MISTRAL_API_KEY not set in environment variables.")

# Initialize Mistral client
client = Mistral(api_key=MISTRAL_API_KEY)
class NamedBytesIO(BytesIO):
    def __init__(self, buffer, name):
        super().__init__(buffer)
        self.name = name

class DeletePDFRequest(BaseModel):
    filename: str


# Authentication with Directus
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")
async def get_current_user(token: str = Depends(oauth2_scheme)):
    try:
        response = requests.get(
            "http://directus:8055/users/me",  # Use 'directus' if in Docker network
            headers={"Authorization": f"Bearer {token}"}
        )
        if response.status_code != 200:
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Invalid authentication credentials",
                headers={"WWW-Authenticate": "Bearer"},
            )
        user_data = response.json().get("data")
        user_data["access_token"] = token
        return user_data
    except Exception as e:
        logger.error(f"Error with token: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Could not validate credentials",
            headers={"WWW-Authenticate": "Bearer"},
        )

# Helper functions (unchanged unless specified)
def download_pdf_from_url(url: str) -> BytesIO:
    """Download a PDF from a URL and return it as a BytesIO object."""
    try:
        response = requests.get(url, stream=True, timeout=10)
        response.raise_for_status()
        if 'application/pdf' not in response.headers.get('Content-Type', ''):
            raise HTTPException(status_code=400, detail=f"URL {url} does not point to a PDF.")
        pdf_bytes = BytesIO(response.content)
        pdf_bytes.name = url.split('/')[-1] or "downloaded_pdf.pdf"
        return pdf_bytes
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error downloading PDF from {url}: {str(e)}")

def get_file_text(files: List[BytesIO]) -> List[Document]:
    """Extract raw text from PDFs using Mistral OCR, and handle other file types as before."""
    documents = []
    for file in files:
        filename = file.name
        if filename.endswith('.pdf'):
            # Reset file pointer to start (important for BytesIO)
            file.seek(0)
            
            # Upload the PDF directly to Mistral for OCR processing
            try:
                uploaded_file = client.files.upload(
                    file={"file_name": file.name, "content": file.read()},
                    purpose="ocr"
                )
                signed_url = client.files.get_signed_url(file_id=uploaded_file.id)
                
                # Process with Mistral OCR
                ocr_response = client.ocr.process(
                    model="mistral-ocr-latest",  # Replace with actual model name if different
                    document={"type": "document_url", "document_url": signed_url.url}
                )
                
                # Extract raw text from OCR response
                full_text = ""
                for page in ocr_response.pages:
                    page_text = page.markdown.strip()  # Assuming markdown output
                    if page_text:
                        full_text += page_text + "\n"  # Concatenate pages with newlines
                
                # Create a single Document with all text
                if full_text:
                    doc = Document(
                        page_content=full_text,
                        metadata={"filename": file.name}
                    )
                    documents.append(doc)
                    # Save extracted text to a .txt file for inspection
                    with open(f"extracted_{file.name}.txt", "w", encoding="utf-8") as text_file:
                        text_file.write(doc.page_content)

            
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Error processing {file.name} with Mistral OCR: {str(e)}")

        # Logic for other file types remains unchanged (simplified without JSON)
        elif filename.endswith('.txt'):
            try:
                text_content = file.read().decode('utf-8').strip()
                if text_content:
                    doc = Document(
                        page_content=text_content,
                        metadata={"filename": file.name}
                    )
                    documents.append(doc)
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Error extracting text from {file.name}: {str(e)}")

        elif filename.endswith('.docx'):
            try:
                temp_path = f"temp_{file.name}"
                with open(temp_path, "wb") as f:
                    f.write(file.read())
                doc = DocxDocument(temp_path)
                text_content = "\n".join(paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip())
                if text_content:
                    doc = Document(
                        page_content=text_content,
                        metadata={"filename": file.name}
                    )
                    documents.append(doc)
                os.remove(temp_path)
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Error extracting text from {file.name}: {str(e)}")

        elif filename.endswith('.ppt') or filename.endswith('.pptx'):
            try:
                temp_path = f"temp_{file.name}"
                with open(temp_path, "wb") as f:
                    f.write(file.read())
                ppt = Presentation(temp_path)
                text_content = ""
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                text_content = text_content.strip()
                if text_content:
                    doc = Document(
                        page_content=text_content,
                        metadata={"filename": file.name}
                    )
                    documents.append(doc)
                os.remove(temp_path)
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"Error extracting text from {file.name}: {str(e)}")

        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.name}")

    return documents




############################# Agent Initialization #############################




# new agent 17/4
# Custom Tool for PDF Content Analysis
class PDFAnalysisTool(BaseTool):
    name: str = "PDF Content Analyzer"
    description: str = "Analyzes text from a random sample of up to five PDF pages using an LLM to classify content type and detect headers for chunking."

    _openai_client = PrivateAttr()
    # func: ClassVar[Callable] = None
    class PDFAnalysisInput(BaseModel):
        text: str = Field(..., description="Text content of the PDF document")

    args_schema: type = PDFAnalysisInput

    # def get_func(self):
    #     return self._run
    
    def __init__(self):
        super().__init__()
        self._openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        object.__setattr__(self, 'func', self._run)

    def _run(self, text: str, filename: str = None) -> Dict[str, Any]:
        """Analyze PDF by sampling up to five pages from extracted text and classifying with LLM."""
        # Sample up to five pages from the extracted text
        sampled_text = self.sample_pages(text)

        # Classify PDF type using LLM
        pdf_type, reasoning = self.classify_pdf_type(sampled_text)

        # Detect headers (simplified heuristic)
        # headers = self.detect_headers(sampled_text)

        # Detect tables (simplified heuristic)
        # tables = self.detect_tables(sampled_text)

        return {
            "pdf_type": pdf_type,
            "text": sampled_text,
            "reasoning": reasoning
        }

    def sample_pages(self, text: str, max_pages: int = 6) -> str:
        """Sample up to five pages from extracted text, assuming page breaks are marked by double newlines."""
        if not text.strip():
            logger.warning("Empty text provided for sampling")
            return ""

        # Split text into pages (Mistral OCR uses markdown with newlines between pages)
        pages = [p.strip() for p in text.split("\n\n") if p.strip()]
        total_pages = len(pages)
        num_pages = min(total_pages, max_pages) if total_pages > 0 else 0

        if num_pages == 0:
            logger.warning("No pages identified in text")
            return text  # Return original text if no pages detected

        # Randomly select pages
        selected_indices = random.sample(range(total_pages), num_pages)
        selected_indices.sort()
        selected_pages = [pages[i] for i in selected_indices]
        logger.info(f"Selected random pages for LLM: {selected_indices}")
        return "\n\n".join(selected_pages)

    def classify_pdf_type(self, text: str) -> Tuple[str, str]:
        """Classify PDF type using OpenAI GPT model."""
        if not text.strip():
            logger.warning("Empty text provided for classification")
            return "normal", "No text available for classification"

        # Truncate text to ~15k words for context window
        max_words = 15000
        words = text.split()
        if len(words) > max_words:
            text = " ".join(words[:max_words])
            logger.info(f"Truncated text to {max_words} words for LLM processing")
        logger.info(f"Text sent to LLM: {text[:500]}... (truncated for log)")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"llm_input_{timestamp}.txt"
        try:
            with open(filename, "w", encoding="utf-8") as f:
                f.write(text)
            logger.info(f"Saved LLM input text to {filename}")
        except Exception as e:
            logger.error(f"Failed to save LLM input text to {filename}: {str(e)}")
        prompt = f"""
        You are an expert in classifying PDF documents based on their content. Given the text extracted from a random sample of up to five pages, classify the PDF as one of the following:
        - "Q&A": Contains question-answer pairs, FAQs, or interrogative sentences followed by responses (e.g., "What is X? Y is…" or "Q: … A: …")  or many question and answers written in numbers index.
        - "table-heavy": Contains tabular data, often with delimiters like `|` or aligned columns.
        - "normal": Primarily narrative or descriptive text with few questions or tables.

        Ignore headers unless they indicate Q&A content (e.g., "FAQ", "Questions"). If headers like "FAQ" or "Questions" appear, prioritize the following text as Q&A.

        Analyze the following text and return a JSON object with:
        - "pdf_type": The classified type ("Q&A", "table-heavy", or "normal").
        - "reasoning": A brief explanation of why you chose this type.
        - "confidence": A confidence score (0.0 to 1.0) for the classification.

        Text:
        {text}

        Output format:
        {{
            "pdf_type": "<type>",
            "reasoning": "<explanation>",
            "confidence": <float>
        }}
        """

        try:
            response = self._openai_client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                temperature=0,
                max_tokens=500,
                response_format={"type": "json_object"}
            )
            result = json.loads(response.choices[0].message.content)
            pdf_type = result.get("pdf_type", "normal")
            reasoning = result.get("reasoning", "No reasoning provided")
            confidence = result.get("confidence", 0.5)

            if pdf_type not in ["Q&A", "table-heavy", "normal"]:
                logger.warning(f"Invalid pdf_type {pdf_type}, defaulting to normal")
                pdf_type = "normal"
                reasoning = f"Invalid type returned, defaulted to normal: {reasoning}"

            logger.info(f"LLM classification: type={pdf_type}, confidence={confidence}, reasoning={reasoning}")
            return pdf_type, reasoning

        except Exception as e:
            logger.error(f"LLM classification failed: {str(e)}")
            return "normal", f"Classification failed due to error: {str(e)}"

    # def detect_headers(self, text: str) -> List[Tuple[str, int]]:
    #     """Simplified header detection for chunking compatibility."""
    #     headers = []
    #     lines = text.splitlines()

    #     for i, line in enumerate(lines):
    #         line = line.strip()
    #         if not line:
    #             continue

    #         words = line.split()
    #         if (len(words) < 10 and (
    #             sum(1 for word in words if word and word[0].isupper()) / max(1, len(words)) > 0.7 or
    #             line.endswith(":")
    #         )):
    #             headers.append((line, text.find(line)))

    #     return headers

    # def detect_tables(self, text: str) -> List[Tuple[str, int, int]]:
    #     """Simplified table detection for chunking compatibility."""
    #     tables = []
    #     md_table_matches = list(re.finditer(r"^\|.*\|$\n?(\|[-:|\s]+\|\n)?((?:\|.*\|\n?)*)", text, re.MULTILINE))
    #     for match in md_table_matches:
    #         tables.append((match.group(0), match.start(), match.end()))

    #     return tables
    
    # def identify_sections(self, text: str, headers: List[tuple], tables: List[tuple]) -> List[Dict[str, Any]]:
    #     """
    #     Identify logical sections in the document for better chunking.
    #     Returns sections with their type, content, and boundaries.
    #     """
    #     sections = []
    #     markers = []
        
    #     # Add headers as section markers
    #     for header_text, pos in headers:
    #         markers.append({"type": "header", "position": pos, "text": header_text})
        
    #     # Add tables as section markers
    #     for table_text, start, end in tables:
    #         markers.append({"type": "table", "position": start, "text": table_text[:50] + "...", "end": end})
        
    #     # Sort all markers by position
    #     markers.sort(key=lambda x: x["position"])
        
    #     # Create sections from markers
    #     text_length = len(text)
    #     start_pos = 0
        
    #     for i, marker in enumerate(markers):
    #         # If there's text before this marker, add it as a section
    #         if marker["position"] > start_pos:
    #             section_text = text[start_pos:marker["position"]].strip()
    #             if section_text:
    #                 sections.append({
    #                     "type": "text",
    #                     "start": start_pos,
    #                     "end": marker["position"],
    #                     "content": section_text[:100] + "..." if len(section_text) > 100 else section_text
    #                 })
            
    #         # Add the marker itself as a section
    #         if marker["type"] == "header":
    #             sections.append({
    #                 "type": "header",
    #                 "start": marker["position"],
    #                 "end": marker["position"] + len(marker["text"]),
    #                 "content": marker["text"]
    #             })
    #             start_pos = marker["position"] + len(marker["text"])
    #         elif marker["type"] == "table":
    #             sections.append({
    #                 "type": "table",
    #                 "start": marker["position"],
    #                 "end": marker["end"],
    #                 "content": marker["text"]
    #             })
    #             start_pos = marker["end"]
        
    #     # Add any remaining text as the final section
    #     if start_pos < text_length:
    #         final_text = text[start_pos:].strip()
    #         if final_text:
    #             sections.append({
    #                 "type": "text",
    #                 "start": start_pos,
    #                 "end": text_length,
    #                 "content": final_text[:100] + "..." if len(final_text) > 100 else final_text
    #             })
        
    #     return sections

# CrewAI Agent for PDF Classification and Chunking
def create_chunking_agent() -> Agent:
    return Agent(
        role="PDF Chunking Specialist",
        goal="Classify PDF type and chunk content intelligently to preserve semantic context.",
        backstory="""You are an expert in document analysis with a deep understanding of PDF structures. 
        You classify PDFs into textbook, Q&A, or table-heavy types and chunk their content to ensure 
        topics, questions, or tables remain intact for accurate retrieval.""",
        tools=[PDFAnalysisTool()],
        verbose=True,
        allow_delegation=False
    )

# Chunking Logic Functions
def chunk_normal(text: str, filename: str, min_chunk_size: int = 400, max_chunk_size: int = 2000) -> List[Document]:
    """
    Chunk a normal PDF text with the following rules:
    - New chunks start at headings (# or ## or ###)
    - Tables should never be split
    - Complete topics should remain together
    - Chunks stay within specified size limits
    
    Args:
        text (str): The OCR'd text from the PDF extracted by Mistral OCR
        filename (str): Name of the PDF file
        min_chunk_size (int): Minimum chunk size in characters (default: 100)
        max_chunk_size (int): Maximum chunk size in characters (default: 1000)
    
    Returns:
        List[Document]: List of Document objects with chunked content
    """
    from langchain.schema import Document
    import re
    
    temp_chunks = []
    current_chunk = ""
    lines = text.split('\n')
    i = 0
    
    # Track if we're inside a table
    in_table = False
    
    while i < len(lines):
        line = lines[i]
        
        # Check if this would exceed max size
        if len(current_chunk + line + "\n") > max_chunk_size and not in_table and current_chunk.strip():
            # Don't split if this is a heading (we'll start a new chunk with it)
            if not re.match(r'^#{1,3}\s+', line.strip()):
                temp_chunks.append(current_chunk.strip())
                current_chunk = ""
        
        # Check if this is a heading (start of a new topic)
        if re.match(r'^#{1,3}\s+', line.strip()) and not in_table:
            # If we already have content in the current chunk, save it before starting a new one
            if current_chunk.strip():
                temp_chunks.append(current_chunk.strip())
                current_chunk = ""
            
            # Add the heading to the new chunk
            current_chunk = line + "\n"
        
        # Check if this is the start of a table
        elif line.strip().startswith('|') or (i+1 < len(lines) and lines[i+1].strip().startswith('| :--')):
            # Check if adding the table would exceed max_chunk_size
            table_content = line + "\n"
            table_end = i
            
            # Find the end of the table to calculate its size
            for j in range(i+1, len(lines)):
                if lines[j].strip().startswith('|'):
                    table_content += lines[j] + "\n"
                    table_end = j
                else:
                    break
            
            # If table itself exceeds max_chunk_size, we have to include it anyway (tables shouldn't be split)
            # If current_chunk + table exceeds max_chunk_size, start a new chunk for the table
            if len(current_chunk) > 0 and len(current_chunk) + len(table_content) > max_chunk_size:
                temp_chunks.append(current_chunk.strip())
                current_chunk = table_content
            else:
                current_chunk += table_content
            
            # Skip ahead to after the table
            i = table_end
            in_table = False
        
        # Regular text
        else:
            current_chunk += line + "\n"
        
        i += 1
    
    # Add the last chunk if it has content
    if current_chunk.strip():
        temp_chunks.append(current_chunk.strip())
    
    # Handle edge case: if no chunks were found, create one chunk from the whole text
    if not temp_chunks:
        # If the text is larger than max_chunk_size, split it into approximately equal parts
        if len(text) > max_chunk_size:
            # Calculate how many chunks we need
            num_chunks = (len(text) + max_chunk_size - 1) // max_chunk_size
            chunk_size = len(text) // num_chunks
            
            for i in range(0, len(text), chunk_size):
                end = min(i + chunk_size, len(text))
                temp_chunks.append(text[i:end].strip())
        else:
            temp_chunks.append(text.strip())
    
    # Process chunks to ensure they meet size requirements
    final_chunks = []
    buffer = ""
    
    for chunk in temp_chunks:
        # If the buffer is empty and the chunk is too small, add to buffer
        if not buffer and len(chunk) < min_chunk_size:
            buffer = chunk
        # If the buffer has content and adding this chunk would still be under max size
        elif buffer and len(buffer) + len(chunk) + 1 <= max_chunk_size:
            buffer = buffer + "\n\n" + chunk
        # If buffer has content but adding chunk would exceed max size
        elif buffer:
            final_chunks.append(Document(
                page_content=buffer,
                metadata={"source": filename}
            ))
            # Start new buffer with current chunk if it's small
            if len(chunk) < min_chunk_size:
                buffer = chunk
            else:
                final_chunks.append(Document(
                    page_content=chunk,
                    metadata={"source": filename}
                ))
                buffer = ""
        # No buffer and chunk is large enough
        else:
            final_chunks.append(Document(
                page_content=chunk,
                metadata={"source": filename}
            ))
    
    # Don't forget the last buffer if it has content
    if buffer:
        final_chunks.append(Document(
            page_content=buffer,
            metadata={"source": filename}
        ))
    
    return final_chunks

def chunk_qa(text: str, filename: str, section: str = None) -> List[Document]:
    """
    Chunk Q&A text by grouping exactly 2 question-answer pairs per chunk.
    Each chunk will contain two complete Q&A pairs, with the last chunk potentially 
    containing just one pair if there's an odd number of pairs.
    """
    # Initialize
    chunks = []
    lines = text.splitlines()
    qa_pairs = []
    current_question = None
    current_answer = []
    min_chunk_size = 10  # Even shorter minimum to catch all Q&A pairs
    
    logger.debug(f"Starting chunk_qa for {filename}, text length: {len(text)}")
    
    try:
        # Based on your sample, most questions appear to be numbered with numbers and parentheses
        # Example: "622) Who is the first Indian to pass the ICS exam?"
        # Example: "623) What is the expansion of IREDA?"
        question_pattern = r"^\s*(\d+[\.\)]+)\s+(.+)$"
        
        # Process each line
        is_first_question = True
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue  # Skip empty lines
            
            # Check if this line matches our question pattern
            question_match = re.match(question_pattern, line)
            
            if question_match:
                # If we already have a question, save the previous Q&A pair
                if current_question and not is_first_question:
                    qa_text = current_question + "\n" + "\n".join(current_answer)
                    if len(qa_text) >= min_chunk_size:
                        qa_pairs.append(qa_text)
                
                # Start a new Q&A pair
                current_question = line
                current_answer = []
                is_first_question = False
            else:
                # Add to the current answer if we have a question
                if current_question:
                    current_answer.append(line)
        
        # Add the last Q&A pair
        if current_question and current_answer:
            qa_text = current_question + "\n" + "\n".join(current_answer)
            if len(qa_text) >= min_chunk_size:
                qa_pairs.append(qa_text)
        
        # Log results
        if not qa_pairs:
            logger.warning(f"No Q&A pairs identified in {filename}. Text sample: {text[:100]}...")
            return []  # Return empty list, not None
        
        logger.info(f"Successfully identified {len(qa_pairs)} Q&A pairs in {filename}")
        
        # Group Q&A pairs into chunks of 2
        for i in range(0, len(qa_pairs), 2):
            if i + 1 < len(qa_pairs):  # We have 2 pairs available
                chunk_content = qa_pairs[i] + "\n\n" + qa_pairs[i+1]
                chunks.append(Document(
                    page_content=chunk_content,
                    metadata={"source": filename, "type": "Q&A", "section": section, "pair_count": 2}
                ))
            else:  # Only 1 pair left
                chunk_content = qa_pairs[i]
                chunks.append(Document(
                    page_content=chunk_content,
                    metadata={"source": filename, "type": "Q&A", "section": section, "pair_count": 1}
                ))
        
    except Exception as e:
        logger.error(f"Error in chunk_qa for {filename}: {str(e)}", exc_info=True)
        return []  # Return empty list on error
    
    logger.info(f"Completed chunk_qa for {filename}, created {len(chunks)} chunks from {len(qa_pairs)} Q&A pairs")
    return chunks

# def chunk_table(text: str, filename: str) -> List[Document]:
#     """
#     Chunk a normal PDF text with the following rules:
#     - New chunks start at headings (# or ## or ###)
#     - Tables should never be split (if they are, table headers are repeated in each chunk)
#     - Maximum chunk size is 3000 characters
#     - Minimum chunk size is 300 characters
#     - Complete topics should remain together when possible
    
#     Args:
#         text (str): The OCR'd text from the PDF extracted by Mistral OCR
#         filename (str): Name of the PDF file
    
#     Returns:
#         List[Document]: List of Document objects with chunked content
#     """
#     from langchain.schema import Document
#     import re
    
#     MAX_CHUNK_SIZE = 3000
#     MIN_CHUNK_SIZE = 300
    
#     chunks = []
#     current_chunk = ""
#     lines = text.split('\n')
#     i = 0
    
#     # Track if we're inside a table
#     in_table = False
#     current_table_headers = []
    
#     # Find the true headers of a table - they typically appear at the start of a table
#     def find_table_headers(start_idx):
#         # Minimum pattern for a table header: text row followed by separator row
#         if start_idx + 1 >= len(lines):
#             return []
            
#         first_line = lines[start_idx].strip()
#         second_line = lines[start_idx + 1].strip()
        
#         # True header has text in first row and separator pattern in second row
#         has_alphabetic = bool(re.search(r'[a-zA-Z]', first_line))
#         is_separator = second_line.startswith('|') and ':--' in second_line
        
#         # If we found a header pattern
#         if has_alphabetic and is_separator:
#             headers = [lines[start_idx], lines[start_idx + 1]]
            
#             # Check for a possible third header row (like Min/Max)
#             if start_idx + 2 < len(lines) and lines[start_idx + 2].strip().startswith('|'):
#                 third_line = lines[start_idx + 2].strip()
#                 # Only include the third line if it has descriptive text
#                 if 'Min' in third_line or 'Max' in third_line or bool(re.search(r'[a-zA-Z]', third_line)):
#                     headers.append(lines[start_idx + 2])
            
#             return headers
#         else:
#             return []
    
#     while i < len(lines):
#         line = lines[i].strip()
        
#         # Check if this is a heading (start of a new topic)
#         if re.match(r'^#{1,3}\s+', line) and not in_table:
#             # If we already have content in the current chunk, save it before starting a new one
#             if current_chunk.strip():
#                 chunks.append(Document(
#                     page_content=current_chunk.strip(),
#                     metadata={"source": filename}
#                 ))
#                 current_chunk = ""
            
#             # Add the heading to the new chunk
#             current_chunk = lines[i] + "\n"
        
#         # Check if this is the start of a new table
#         elif line.startswith('|') and not in_table:
#             # Try to identify if this is a true table header
#             headers = find_table_headers(i)
            
#             if headers:
#                 in_table = True
#                 current_table_headers = headers
                
#                 # Add all header rows to current chunk
#                 for header in headers:
#                     current_chunk += header + "\n"
                
#                 # Skip ahead past the header rows we've already processed
#                 i += len(headers) - 1
#             else:
#                 # Not a header, just a regular line with | characters
#                 current_chunk += lines[i] + "\n"
        
#         # If we're in a table and this line has pipe characters
#         elif in_table and line.startswith('|'):
#             # Check if adding this line would exceed the max size
#             if len(current_chunk) + len(line) > MAX_CHUNK_SIZE:
#                 # Save current chunk
#                 chunks.append(Document(
#                     page_content=current_chunk.strip(),
#                     metadata={"source": filename}
#                 ))
                
#                 # Start a new chunk with the table headers
#                 current_chunk = ""
#                 for header in current_table_headers:
#                     current_chunk += header + "\n"
                
#                 # Now add the current line
#                 current_chunk += lines[i] + "\n"
#             else:
#                 # Add the line to the current chunk
#                 current_chunk += lines[i] + "\n"
            
#             # Check if next line doesn't have pipe characters, indicating end of table
#             if i+1 >= len(lines) or not lines[i+1].strip().startswith('|'):
#                 in_table = False
        
#         # Regular text
#         else:
#             # End of table detection
#             if in_table:
#                 in_table = False
            
#             # Check if adding this line would exceed max size
#             if len(current_chunk) + len(lines[i]) > MAX_CHUNK_SIZE and len(current_chunk) >= MIN_CHUNK_SIZE:
#                 # Save current chunk
#                 chunks.append(Document(
#                     page_content=current_chunk.strip(),
#                     metadata={"source": filename}
#                 ))
#                 current_chunk = lines[i] + "\n"
#             else:
#                 current_chunk += lines[i] + "\n"
        
#         i += 1
    
#     # Add the last chunk if it has content
#     if current_chunk.strip():
#         chunks.append(Document(
#             page_content=current_chunk.strip(),
#             metadata={"source": filename}
#         ))
    
#     # Post-processing: merge small chunks with previous ones if possible
#     processed_chunks = []
#     for i, chunk in enumerate(chunks):
#         # If this is a small chunk and not the first one, try to merge with previous
#         if len(chunk.page_content) < MIN_CHUNK_SIZE and i > 0:
#             prev_chunk = processed_chunks[-1]
#             if len(prev_chunk.page_content) + len(chunk.page_content) <= MAX_CHUNK_SIZE:
#                 # Merge with previous chunk
#                 processed_chunks[-1] = Document(
#                     page_content=prev_chunk.page_content + "\n\n" + chunk.page_content,
#                     metadata=prev_chunk.metadata
#                 )
#             else:
#                 # Can't merge, keep as is
#                 processed_chunks.append(chunk)
#         else:
#             processed_chunks.append(chunk)
    
#     # Handle edge case: if no chunks were created, create one from the whole text
#     if not processed_chunks:
#         processed_chunks.append(Document(
#             page_content=text.strip(),
#             metadata={"source": filename}
#         ))
    
#     return processed_chunks

def chunk_table(text: str, filename: str) -> List[Document]:
    """
    Memory-efficient version of the chunking function
    - New chunks start at headings (# or ## or ###)
    - Tables with headers always have headers repeated when split
    - Maximum chunk size is 3000 characters
    - Minimum chunk size is 300 characters
    
    Args:
        text (str): The OCR'd text from the PDF
        filename (str): Name of the PDF file
    
    Returns:
        List[Document]: List of Document objects with chunked content
    """
    from langchain.schema import Document
    import re
    
    MAX_CHUNK_SIZE = 2000
    MIN_CHUNK_SIZE = 300
    
    # Process text line by line instead of storing all lines in memory
    lines = text.split('\n')
    
    chunks = []
    current_chunk = ""
    i = 0
    
    # Track table state with minimal memory usage
    in_table = False
    current_table_headers = []
    
    # Simplified header detection
    def is_table_header_row(idx):
        if idx + 1 >= len(lines):
            return False
        first_line = lines[idx].strip()
        second_line = lines[idx + 1].strip()
        return (first_line.startswith('|') and 
                bool(re.search(r'[a-zA-Z]', first_line)) and 
                second_line.startswith('|') and 
                (':--' in second_line or '---' in second_line))
    
    while i < len(lines):
        line = lines[i].strip()
        
        # Check if this is a heading
        if re.match(r'^#{1,3}\s+', line) and not in_table:
            if current_chunk.strip():
                chunks.append(Document(
                    page_content=current_chunk.strip(),
                    metadata={"source": filename}
                ))
                current_chunk = ""
            current_chunk = lines[i] + "\n"
        
        # Check if this is the start of a table
        elif line.startswith('|') and not in_table and is_table_header_row(i):
            in_table = True
            # Store headers (reuse list object to save memory)
            current_table_headers = []
            
            # Add header and separator row
            current_table_headers.append(lines[i])
            current_table_headers.append(lines[i+1])
            current_chunk += lines[i] + "\n" + lines[i+1] + "\n"
            
            # Check for Min/Max subheader
            if i+2 < len(lines) and lines[i+2].strip().startswith('|'):
                third_line = lines[i+2].strip()
                if 'Min' in third_line or 'Max' in third_line:
                    current_table_headers.append(lines[i+2])
                    current_chunk += lines[i+2] + "\n"
                    i += 2  # Skip two lines
                else:
                    i += 1  # Skip one line
            else:
                i += 1  # Skip one line
            
            continue
        
        # Process table data rows
        elif in_table and line.startswith('|'):
            # Check if we need to start a new chunk
            if len(current_chunk) + len(line) + 10 > MAX_CHUNK_SIZE:  # Add buffer
                chunks.append(Document(
                    page_content=current_chunk.strip(),
                    metadata={"source": filename}
                ))
                # Start new chunk with headers
                current_chunk = "\n".join(current_table_headers) + "\n"
            
            current_chunk += lines[i] + "\n"
            
            # Check if next line ends the table
            if i+1 >= len(lines) or not lines[i+1].strip().startswith('|'):
                in_table = False
        
        # Regular text
        else:
            if in_table:
                in_table = False
            
            if len(current_chunk) + len(line) + 10 > MAX_CHUNK_SIZE and len(current_chunk) >= MIN_CHUNK_SIZE:
                chunks.append(Document(
                    page_content=current_chunk.strip(),
                    metadata={"source": filename}
                ))
                current_chunk = lines[i] + "\n"
            else:
                current_chunk += lines[i] + "\n"
        
        i += 1
    
    # Add the final chunk
    if current_chunk.strip():
        chunks.append(Document(
            page_content=current_chunk.strip(),
            metadata={"source": filename}
        ))
    
    # Simplified post-processing for small chunks
    i = 0
    while i < len(chunks) - 1:
        if len(chunks[i].page_content) < MIN_CHUNK_SIZE:
            # Try to merge with next chunk
            if len(chunks[i].page_content) + len(chunks[i+1].page_content) <= MAX_CHUNK_SIZE:
                chunks[i+1] = Document(
                    page_content=chunks[i].page_content + "\n\n" + chunks[i+1].page_content,
                    metadata=chunks[i+1].metadata
                )
                chunks.pop(i)
                continue
        i += 1
    
    return chunks


# # 17/4 with agent
# # Modified get_text_chunks Function
# def get_text_chunks(documents: List[Document], user_id: str) -> List[Document]:
#     """Chunk documents using a CrewAI agent and add to ChromaDB."""
#     from db import add_to_collection, overwriting_chromadb
    
#     logger.info(f"Starting get_text_chunks with {len(documents)} documents for user {user_id}")
#     all_chunks = []
#     agent = create_chunking_agent()
#     analysis_tool = PDFAnalysisTool()  
    
#     for doc in documents:
#         try:
#             filename = doc.metadata["filename"]
#             logger.info(f"Processing document: {filename}, text length: {len(doc.page_content)}")
            
#             full_text = doc.page_content 
#             # print("full text.......", full_text[50])
#             logger.debug(f"Calling overwriting_chromadb for {filename}")
#             overwriting_chromadb(user_id, filename)
#             logger.debug(f"Completed overwriting_chromadb for {filename}")
            
#             logger.debug(f"Running PDFAnalysisTool on {filename}")
#             analysis = analysis_tool._run(doc.page_content, filename=filename)
#             logger.info(f"Analysis result for {filename}: pdf_type={analysis['pdf_type']}, "
#                        f"reasoning={analysis['reasoning']}")
            
#             pdf_type = analysis["pdf_type"]
#             # headers = analysis["headers"]
#             text = analysis["text"]
#             # print('text:...........', text)
#             if pdf_type == "Q&A":
#                 logger.info(f"Processing Q&A document {filename} as a whole")
#                 chunks = chunk_qa(full_text, filename)
#                 if not chunks:
#                     logger.warning(f"Q&A chunking failed for {filename}, falling back to standard chunking")
#                     chunks = chunk_normal(full_text, filename)
#             elif pdf_type=="table-heavy":
#                 logger.info(f'processing table document{filename} table heavy')
#                 chunks=chunk_table(full_text, filename)
#                 if not chunks:
#                     logger.warning(f"table chunking failed for {filename}, falling back to standard chunking")
#                     # chunks = chunk_normal(full_text, filename)
#             else:
#                 logger.info(f"Processing normal document {filename} as a whole")
#                 chunks = chunk_normal(full_text, filename)
                
#                 # chunks = []
#                 # last_pos = 0
#                 # current_section = None
                
#                 # for header, pos in headers:
#                 #     if last_pos < pos:
#                 #         section_text = full_text[last_pos:pos].strip()
#                 #         current_section = header
#                 #         if section_text:
#                 #             if pdf_type == "table-heavy":
#                 #                 new_chunks = chunk_table(section_text, filename, current_section)
#                 #             else:
#                 #                 new_chunks = chunk_normal(section_text, filename, current_section)
#                 #             chunks.extend(new_chunks)
#                 #     last_pos = pos + len(header)
                
#                 # if last_pos < len(text):
#                 #     section_text = text[last_pos:].strip()
#                 #     if section_text:
#                 #         if pdf_type == "table-heavy":
#                 #             new_chunks = chunk_table(section_text, filename, current_section)
#                 #         else:
#                 #             new_chunks = chunk_normal(section_text, filename, current_section)
#                 #         chunks.extend(new_chunks)
            
#             if not chunks:
#                 logger.warning(f"No chunks created for {filename}")
#                 continue
            
#             logger.info(f"Created {len(chunks)} chunks for {filename}")
            
#             chunk_texts = [chunk.page_content for chunk in chunks]
#             logger.info("moving the saving on chromadb")
#             add_to_collection(chunk_texts, filename, user_id)
            
#             all_chunks.extend(chunks)
            
#         except Exception as e:
#             logger.error(f"Error processing {doc.metadata['filename']}: {str(e)}", exc_info=True)
#             raise HTTPException(status_code=500, detail=f"Error processing {doc.metadata['filename']}: {str(e)}")
    
#     logger.info(f"Completed get_text_chunks, returning {len(all_chunks)} chunks")
#     return all_chunks

import re
import json
from typing import List, Dict, Any, Optional
from langchain.schema.document import Document
import logging

def extract_tables_from_markdown(text: str) -> List[Dict[str, Any]]:
    """Extracts markdown tables from text and returns list of {table_text, start_index, end_index}"""
    # Look for a header line with | characters
    # Followed by a separator line with |, :, and - characters
    # Followed by one or more data lines with | characters
    table_pattern = r'((?:(?:\#[^\n]*\n)?)?(?:\|[^\n]*\|\n)+\|[\s:-]*\|[\s:-]*\|.*?(?=\n\n|\n[^|]|\Z))'
    tables = []
   
    for match in re.finditer(table_pattern, text, re.DOTALL):
        # Check if this is actually a table by looking for pipe and separator characters
        table_text = match.group(0)
        lines = table_text.strip().split('\n')
       
        # Verify this is a valid markdown table (at least 2 lines, with proper separator)
        if len(lines) >= 2 and all('|' in line for line in lines[:2]):
            # Look for a preceding title (# Header)
            start_index = match.start()
            preceding_text = text[:start_index]
            title_match = re.search(r'((?:^|\n)#\s+[^\n]+)\n*$', preceding_text)
            
            if title_match:
                # Include the title in the table text
                table_text = title_match.group(1) + '\n' + table_text
                start_index = title_match.start()
            
            tables.append({
                "table_text": table_text,
                "start_index": start_index,
                "end_index": match.end()
            })
    
    return tables

def parse_markdown_table(table_text: str) -> Dict[str, Any]:
    """Parse markdown table into structured JSON format with dynamic structure based on headers"""
    lines = table_text.strip().split('\n')
    
    # Check if there are enough lines to form a table
    if len(lines) < 2:
        return {"content": table_text, "error": "Not enough lines to form a table"}
    
    # Extract header (assuming first line is header)
    headers = [h.strip() for h in lines[0].strip('|').split('|')]
    headers = [h for h in headers if h]  # Remove empty headers
    
    # Extract title if it exists (assuming it's before the table)
    title = ""
    title_match = re.search(r'#\s+([^\n]+)', table_text)
    if title_match:
        title = title_match.group(1).strip()
    
    # Skip separator line (second line with dashes and colons)
    start_idx = 1
    while start_idx < len(lines) and ('---' in lines[start_idx] or ':--:' in lines[start_idx] or '-:' in lines[start_idx] or ':-' in lines[start_idx]):
        start_idx += 1
    
    data_lines = lines[start_idx:] if start_idx < len(lines) else []
    
    # Build a dynamic table structure
    table_data = []
    
    # Track row continuation for parameters that span multiple lines
    current_row = None
    current_id_col = 0  # Assume first column is the ID column by default
    
    # Find which column might be an ID column (this is just a heuristic)
    for i, header in enumerate(headers):
        if 'id' in header.lower() or 'parameter' in header.lower() and 'id' in header.lower():
            current_id_col = i
            break
    
    for line in data_lines:
        cells = [cell.strip() for cell in line.strip('|').split('|')]
        cells = [c for c in cells if c != '']  # Remove empty cells from split
        
        # Skip completely empty lines or separator lines
        if not cells or all(cell == '' or cell.startswith(':--:') or cell == '-' for cell in cells):
            continue
            
        # Ensure cells and headers have same length for mapping
        while len(cells) < len(headers):
            cells.append('')
            
        # If cells has more entries than headers, truncate
        cells = cells[:len(headers)] if len(cells) > len(headers) else cells
        
        # Check if this is a continuation row or a new row
        is_new_row = len(cells) > current_id_col and cells[current_id_col] != ''
        
        if is_new_row:
            # If we have a current row, add it to our table data
            if current_row:
                table_data.append(current_row)
            
            # Create a new row with all available data
            current_row = {}
            for i, header in enumerate(headers):
                if i < len(cells) and cells[i]:
                    current_row[header] = cells[i]
        elif current_row:
            # This is a continuation row, update existing values or add new ones
            for i, header in enumerate(headers):
                if i < len(cells) and cells[i]:
                    # If the header already exists in current_row, append the new value
                    if header in current_row:
                        # Only append if it's not a duplicate
                        current_value = current_row[header]
                        if isinstance(current_value, list):
                            if cells[i] not in current_value:
                                current_row[header].append(cells[i])
                        else:
                            if cells[i] != current_value:
                                current_row[header] = [current_value, cells[i]]
                    else:
                        current_row[header] = cells[i]
    
    # Add the last row if it exists
    if current_row:
        table_data.append(current_row)
    
    # Build the final result
    result = {
        "table_type": "dynamic",
        "headers": headers,
        "rows": table_data
    }
    
    if title:
        result["title"] = title
        
    return result

def convert_markdown_tables_to_json(text: str) -> str:
    """Converts markdown tables in text to JSON format and returns the modified text"""
    tables = extract_tables_from_markdown(text)
    
    if not tables:
        return text
    
    # Process tables in reverse order to maintain correct indices
    tables.sort(key=lambda x: x["start_index"], reverse=True)
    modified_text = text
    
    for table in tables:
        try:
            # Get table section
            table_text = table["table_text"]
            
            # Check if this is actually a table (at least 2 lines with pipe characters)
            table_lines = table_text.strip().split('\n')
            if len(table_lines) < 2 or not all('|' in line for line in table_lines[:2]):
                continue
                
            # Parse table to JSON structure
            json_data = parse_markdown_table(table_text)
            
            # Convert to pretty JSON string
            json_text = json.dumps(json_data, indent=2)
            
            # Replace table with JSON in the text
            modified_text = (
                modified_text[:table["start_index"]] +
                json_text +
                modified_text[table["end_index"]:]
            )
        except Exception as e:
            logger.warning(f"Error converting markdown table to JSON: {str(e)}")
            # If there's an error, keep the original table
            continue
    
    return modified_text

def get_text_chunks(documents: List[Document], user_id: str) -> List[Document]:
    """Chunk documents using a CrewAI agent and add to ChromaDB."""
    from db import add_to_collection, overwriting_chromadb
    
    logger.info(f"Starting get_text_chunks with {len(documents)} documents for user {user_id}")
    all_chunks = []
    agent = create_chunking_agent()
    analysis_tool = PDFAnalysisTool()  
    
    for doc in documents:
        try:
            filename = doc.metadata["filename"]
            logger.info(f"Processing document: {filename}, text length: {len(doc.page_content)}")
            
            full_text = doc.page_content 
            logger.debug(f"Calling overwriting_chromadb for {filename}")
            overwriting_chromadb(user_id, filename)
            logger.debug(f"Completed overwriting_chromadb for {filename}")
            
            logger.debug(f"Running PDFAnalysisTool on {filename}")
            analysis = analysis_tool._run(doc.page_content, filename=filename)
            logger.info(f"Analysis result for {filename}: pdf_type={analysis['pdf_type']}, "
                       f"reasoning={analysis['reasoning']}")
            
            pdf_type = analysis["pdf_type"]
            text = analysis["text"]
            
            if pdf_type == "Q&A":
                logger.info(f"Processing Q&A document {filename} as a whole")
                chunks = chunk_qa(full_text, filename)
                if not chunks:
                    logger.warning(f"Q&A chunking failed for {filename}, falling back to standard chunking")
                    chunks = chunk_normal(full_text, filename)
            elif pdf_type == "table-heavy":
                logger.info(f'Processing table document {filename} table heavy')
                chunks = chunk_table(full_text, filename)
                if not chunks:
                    logger.warning(f"Table chunking failed for {filename}, falling back to standard chunking")
                    # chunks = chunk_normal(full_text, filename)
            else:
                logger.info(f"Processing normal document {filename} as a whole")
                chunks = chunk_normal(full_text, filename)
            
            if not chunks:
                logger.warning(f"No chunks created for {filename}")
                continue
            
            logger.info(f"Created {len(chunks)} chunks for {filename}")
            
            # Process chunks to convert markdown tables to JSON
            processed_chunks = []
            for chunk in chunks:
                # Convert markdown tables in the chunk to JSON format
                processed_content = convert_markdown_tables_to_json(chunk.page_content)
                
                # Create a new document with processed content but keep same metadata
                processed_chunk = Document(
                    page_content=processed_content,
                    metadata=chunk.metadata
                )
                processed_chunks.append(processed_chunk)
            
            # Use the processed chunks for storage and return
            chunk_texts = [chunk.page_content for chunk in processed_chunks]
            logger.info("Saving the JSON-formatted chunks to ChromaDB")
            add_to_collection(chunk_texts, filename, user_id)
            
            all_chunks.extend(processed_chunks)
            
        except Exception as e:
            logger.error(f"Error processing {doc.metadata['filename']}: {str(e)}", exc_info=True)
            raise HTTPException(status_code=500, detail=f"Error processing {doc.metadata['filename']}: {str(e)}")
    
    logger.info(f"Completed get_text_chunks, returning {len(all_chunks)} chunks")
    return all_chunks

def format_docs(docs, sources):
    return "\n\n".join(f"Source: {source}\n{doc_text}" for doc_text, source in zip(docs, sources))

def save_to_directus(query: str, response: str, conversation_id: int = None):
    if not isinstance(response, str):
        response = str(response)
    if conversation_id is None:
        try:
            response_get = requests.get(f"{DIRECTUS_URL}/items/conversation_history?sort=-conversation_id&limit=1", headers=HEADERS, timeout=5)
            response_get.raise_for_status()
            data = response_get.json().get("data", [])
            last_conversation_id = max((entry["conversation_id"] or 0) for entry in data)
            conversation_id = last_conversation_id + 1
        except (requests.exceptions.RequestException, ValueError):
            conversation_id = 1

    payload = {"query": query, "response": response, "conversation_id": conversation_id}
    try:
        response_post = requests.post(f"{DIRECTUS_URL}/items/conversation_history", json=payload, headers=HEADERS)
        response_post.raise_for_status()
    except requests.exceptions.RequestException as e:
        raise HTTPException(status_code=500, detail=f"Error saving to Directus: {str(e)}")

def get_conversation_history(user_id: str):
    """Retrieve the last 5 conversations for a specific user from Directus, sorted by latest conversation_id."""
    try:
        response = requests.get(
            f"{DIRECTUS_URL}/items/conversation_history",
            params={
                "filter[user_created][_eq]": user_id,  # Filter by user
                "sort": "-conversation_id",  # Newest first
                "limit": 5,  # Last 5 conversations
            },
            headers=HEADERS,
            timeout=5
        )
        response.raise_for_status()

        data = response.json().get("data", [])
        # print(f"Raw API response for user {user_id}: {response.text}")  # Log the full response
        print(f"Retrieved history for user {user_id} (order): {[entry['conversation_id'] for entry in data]}")  # Log conversation_ids in order
        return data
    except requests.exceptions.RequestException as e:
        print(f"Warning: Could not connect to Directus for user {user_id}. Proceeding without history: {str(e)}")
        return []

def process_query(user_question: str, user_id) -> str:
    
    logger.info("before retrival")
    # id='da051188-f956-4398-b922-c0fe05a3dd6e' 
    docs = retrieve_from_collection(user_question,user_id,  top_k=8)
    if not docs:
        return "No relevant information found in the provided PDFs."
    logger.info("retrival succesfull-----------------------------")
    sources = [doc['metadata'].get('filename', 'Unknown') for doc in docs]
    doc_texts = [doc['document'] for doc in docs]
    context = format_docs(doc_texts, sources)
    logger.info("before conversation history")
    history = get_conversation_history(user_id)
    history_str = "Previous Conversation History:\n" + "\n".join(
        f"User: {entry['query']}\nAssistant: {entry['response']}" for entry in history
    ) if history else "No previous history."
    logger.info("after conversation history")
    openai_api_key = os.getenv("OPENAI_API_KEY")
    if not openai_api_key:
        raise HTTPException(status_code=500, detail="OPENAI_API_KEY environment variable not set.")
    logger.info("before asking model")
    model = ChatOpenAI(model="gpt-4o", temperature=0.2)
    messages = [
    {
        "role": "system",
            "content": f"""You are a friendly AI assistant that helps users understand PDF documents and engage in general conversation. 
            The following information comes from these PDF files: {', '.join(set(sources))}
            - If the answer is not available in the context, just strictly say: **"Answer is not available in the context."**, don't write anything else.
            **Conversation Memory:** Use this history for context:
            {history_str}
            if the query is anything related to summary or thier synonyms like sum up, etc, ask for pdf name.
            **Behavior:**
            - If the user greets you (e.g., says "hello", "hi", "hey"), respond with a friendly greeting like "Hi, how can I assist you?" or "Hello! How can I help you today?"
            - If the user asks a question and relevant PDF content is available, respond with: **"Sources: [list of PDF filenames, comma-separated]"**  followed by a detailed and structured answer based on the PDF content.
            - If the user asks a question but no relevant PDF content is available, respond with: **"No relevant information found in the provided PDFs."** and offer to help with something else (e.g., "Can I assist you with something else?").
            **Rules:**  
            - Use the provided context and history to answer questions about PDFs.    
            - strictly answer from the context only.
            - always provide the sources in the response
            - Perform calculations if needed.  
            - Ask for clarification if the user's query is unclear.
            - make it structured and organised with proper formating like headings, bullets , new lines, bold characters, etc.
            - Be friendly and conversational in your responses."""
    },
        {"role": "user", "content": f"Question: {user_question}\n\nInformation from PDFs:\n{context}"}
    ]
    response = model.invoke(messages).content
    print("response:  ", response)
    logger.info("after getting response from model")
    save_to_directus(user_question, response)
    logger.info("saving to directus")
    return response

def get_user_conversation_history(user_id: str, limit: int = 50, offset: int = 0):
    try:
        response = requests.get(
            f"{DIRECTUS_URL}/items/conversation_history",
            params={
                "filter[user_created][_eq]": user_id,
                # "sort": "conversation_id",  # Oldest first
                "limit": limit,
                "offset": offset,
            },
            headers=HEADERS,
            timeout=5
        )
        response.raise_for_status()
        data = response.json().get("data", [])
        print(f"Retrieved user history for {user_id}: {[entry['conversation_id'] for entry in data]}")
        return data
    except requests.exceptions.RequestException as e:
        print(f"Error fetching user conversation history: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error fetching conversation history: {str(e)}")

def save_pdf_metadata_to_directus(user_id: str, s3_key: str, original_filename: str, headers: dict) -> dict:
    """
    Save PDF metadata to the Directus 'uploaded_pdfs' collection.
    
    Args:
        user_id (str): The ID of the user who uploaded the PDF.
        s3_key (str): The S3 key (path) where the PDF is stored.
        original_filename (str): The original name of the uploaded file.
        headers (dict): HTTP headers with the authorization token for Directus.
    
    Returns:
        dict: The response data from Directus.
    
    Raises:
        HTTPException: If the save operation fails.
    """
    payload = {
        "user_id": user_id,
        "s3_key": s3_key,
        "original_filename": original_filename,
        "upload_date": datetime.utcnow().isoformat()
    }
    
    try:
        response = requests.post(
            f"{DIRECTUS_URL}/items/uploaded_pdfs",
            json=payload,
            headers=headers,
            timeout=5  # Add a timeout to avoid hanging
        )
        response.raise_for_status()
        logger.info(f"Saved metadata for {original_filename} to Directus: {response.json()}")
        return response.json().get("data", {})
    except requests.exceptions.RequestException as e:
        logger.error(f"Failed to save metadata for {original_filename} to Directus: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Failed to save metadata to Directus: {str(e)}")

# @app.get("/list-pdfs")
# async def list_pdfs(user: dict = Depends(get_current_user)):
#     user_id = user.get("id")
#     if not user_id:
#         raise HTTPException(status_code=400, detail="User ID not found in token.")
    
#     headers = {
#         "Content-Type": "application/json",
#         "Authorization": f"Bearer {user.get('access_token')}"
#     }
    
#     try:
#         response = requests.get(
#             f"{DIRECTUS_URL}/items/uploaded_pdfs",
#             params={
#                 "filter[user_id][_eq]": user_id,
#                 "fields": "original_filename,s3_key",
#                 "sort": "-date_created"
#             }, 
#             headers=headers,
#             timeout=5
#         )
#         response.raise_for_status()
#         pdfs = response.json().get("data", [])
#         # Add presigned URLs for downloading
#         for pdf in pdfs:
#             pdf["download_url"] = s3_client.generate_presigned_url(
#                 "get_object",
#                 Params={"Bucket": S3_BUCKET_NAME, "Key": pdf["s3_key"]},
#                 ExpiresIn=3600  # URL valid for 1 hour
#             )
#         return {"pdfs": pdfs}
#     except requests.exceptions.RequestException as e:
#         logger.error(f"Error fetching PDF list for user {user_id}: {str(e)}")
#         raise HTTPException(status_code=500, detail=f"Error fetching PDF list: {str(e)}")
    

@app.post("/upload-file")
async def upload_file(files: List[UploadFile], user: dict = Depends(get_current_user)):
    if not files:
        raise HTTPException(status_code=400, detail="No files uploaded.")
    
    file_bytes = []
    user_id = user.get("id")
    if not user_id:
        raise HTTPException(status_code=400, detail="User ID not found in token.")
    
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {user.get('access_token')}"
    }
    
    try:
        for file in files:
            if not (file.filename.lower().endswith('.pdf') or file.filename.lower().endswith('.txt') or 
                    file.filename.lower().endswith('.docx') or file.filename.lower().endswith('.ppt') or 
                    file.filename.lower().endswith('.pptx')):
                raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.filename}")
            file_content = await file.read()
            file_bytes.append(NamedBytesIO(file_content, file.filename))
            s3_key = f"{user_id}/{file.filename}"
            logger.info(f"Uploading {file.filename} to S3 bucket {S3_BUCKET_NAME} at {s3_key}")
            try:
                content_type = {
                    '.pdf': 'application/pdf',
                    '.txt': 'text/plain',
                    '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    '.ppt': 'application/vnd.ms-powerpoint',
                    '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                }[os.path.splitext(file.filename.lower())[1]]
                # s3_client.upload_fileobj(
                #     BytesIO(file_content),
                #     S3_BUCKET_NAME,
                #     s3_key,
                #     ExtraArgs={"ContentType": content_type}
                # )
                # logger.info(f"Successfully uploaded {file.filename} to S3 at {s3_key}")
                # save_pdf_metadata_to_directus(user_id, s3_key, file.filename, headers)
            except Exception as s3_error:
                logger.error(f"Failed to upload {file.filename} to S3: {str(s3_error)}")
                raise HTTPException(status_code=500, detail=f"S3 upload failed for {file.filename}: {str(s3_error)}")
        
        documents = get_file_text(file_bytes)
        get_text_chunks(documents, user_id)
        return {"message": f"Processed {len(documents)} files successfully and uploaded to S3."}
    except Exception as e:
        logger.error(f"Error processing files or uploading to S3: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error processing files or uploading to S3: {str(e)}")
    
@app.post("/process-url")
async def process_url(request: URLRequest, user: dict = Depends(get_current_user)):
    if not request.urls:
        raise HTTPException(status_code=400, detail="No URLs provided.")
    user_id = user.get("id")
    if not user_id:
        raise HTTPException(status_code=400, detail="User ID not found in token.")
    
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {user.get('access_token')}"
    }
    
    file_bytes_raw = []
    file_bytes = []
    filenames = []
    
    for url in request.urls:
        response = requests.get(url, stream=True, timeout=10)
        response.raise_for_status()
        filename = url.split("/")[-1].lower()
        if filename.endswith('.pdf'):
            if 'application/pdf' not in response.headers.get('Content-Type', ''):
                raise HTTPException(status_code=400, detail=f"URL {url} does not point to a PDF.")
            raw_content = response.content
            file_bytes_raw.append(raw_content)
            file_bytes.append(NamedBytesIO(raw_content, filename))
            filenames.append(filename)
        elif filename.endswith('.txt'):
            if 'text/plain' not in response.headers.get('Content-Type', ''):
                raise HTTPException(status_code=400, detail=f"URL {url} does not point to a TXT file.")
            raw_content = response.content
            file_bytes_raw.append(raw_content)
            file_bytes.append(NamedBytesIO(raw_content, filename))
            filenames.append(filename)
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type from URL: {url}")

    if not file_bytes:
        raise HTTPException(status_code=400, detail="No valid files downloaded from URLs.")
    
    try:
        for i, file in enumerate(file_bytes):
            filename = filenames[i]
            s3_key = f"{user_id}/{filename}"
            logger.info(f"Uploading {filename} to S3 bucket {S3_BUCKET_NAME} at {s3_key}")
            s3_client.upload_fileobj(
                BytesIO(file_bytes_raw[i]),
                S3_BUCKET_NAME,
                s3_key,
                ExtraArgs={"ContentType": "application/pdf" if filename.endswith('.pdf') else "text/plain"}
            )
            logger.info(f"Successfully uploaded {filename} to S3 at {s3_key}")
            save_pdf_metadata_to_directus(user_id, s3_key, filename, headers)
        
        documents = get_file_text(file_bytes)
        get_text_chunks(documents, user_id)
        return {"message": f"Processed {len(documents)} files from URLs successfully and uploaded to S3."}
    except Exception as e:
        logger.error(f"Error processing files from URLs or uploading to S3: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error processing files from URLs or uploading to S3: {str(e)}")

@app.post("/query", response_model=QueryResponse)
async def query(request: QueryRequest,user: dict = Depends(get_current_user)):
    """Answer a user question based on stored PDFs with conversation history."""
    user_id = user.get("id")
    if not user_id:
        raise HTTPException(status_code=400, detail="User ID not found in token.")
    
    headers = {
        "Content-Type": "application/json",
        "Authorization": f"Bearer {user.get('access_token')}"
    }
    if not request.question:
        raise HTTPException(status_code=400, detail="No question provided.")
    
    try:
        response = process_query(request.question, user_id)
        return {"response": response}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing query: {str(e)}")

@app.get("/conversation-history")
async def fetch_conversation_history(limit: int = 50, offset: int = 0, user: dict = Depends(get_current_user)):
    try:
        global HEADERS
        HEADERS = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {user.get('access_token', 'default-token')}"
        }
        user_id = user.get("id")
        logger.info(f"Fetching conversation history for user {user_id}")
        if not user_id:
            logger.error("User ID not found in token")
            raise HTTPException(status_code=400, detail="User ID not found in token.")
        history = get_user_conversation_history(user_id, limit, offset)
        logger.info(f"Retrieved {len(history)} history items for user {user_id}")
        return {"history": history}
    except Exception as e:
        logger.error(f"Error in fetch_conversation_history: {str(e)}\n{traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"Error fetching conversation history: {str(e)}")

@app.delete("/delete-pdf")
async def delete_pdf(s3_key: str, filename: str, user: dict = Depends(get_current_user)):
    user_id = user.get("id")
    if not user_id:
        raise HTTPException(status_code=400, detail="User ID not found in token.")

    try:
        # Step 1: Delete from S3
        logger.info(f"Deleting {s3_key} from S3 for user {user_id}")
        s3_client.delete_object(Bucket=S3_BUCKET_NAME, Key=s3_key)
        logger.info(f"Successfully deleted {s3_key} from S3")

        # Step 2: Delete from ChromaDB
        delete_pdf_from_collection(user_id, filename)
        logger.info(f"Successfully deleted all data for {filename} from ChromaDB for user {user_id}")

        return {"message": f"PDF {filename} deleted successfully from S3 and ChromaDB."}
    except Exception as e:
        logger.error(f"Error deleting PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error deleting PDF: {str(e)}")
    
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)